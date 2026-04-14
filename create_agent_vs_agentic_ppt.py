#!/usr/bin/env python3
"""
Agent AI vs Agentic AI - Professional PPT Generator
Creates a comprehensive comparison presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme - Modern tech theme
COLORS = {
    'primary': RGBColor(0x1A, 0x73, 0xE8),      # Google Blue
    'secondary': RGBColor(0x34, 0xA8, 0x53),     # Green
    'accent': RGBColor(0xFF, 0x6D, 0x00),        # Orange
    'dark': RGBColor(0x20, 0x21, 0x24),          # Dark gray
    'darker': RGBColor(0x0D, 0x11, 0x17),        # Very dark
    'light': RGBColor(0xF8, 0xF9, 0xFA),         # Light gray
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'gradient_blue': RGBColor(0x42, 0x85, 0xF4),
    'gradient_purple': RGBColor(0x7C, 0x4D, 0xFF),
    'gradient_teal': RGBColor(0x00, 0xBC, 0xD4),
    'text_gray': RGBColor(0x5F, 0x63, 0x68),
    'card_bg': RGBColor(0x2D, 0x2F, 0x31),
    'vs_red': RGBColor(0xFF, 0x44, 0x44),
    'vs_orange': RGBColor(0xFF, 0x8C, 0x00),
}

def set_slide_bg(slide, color):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape_type=MSO_SHAPE.RECTANGLE):
    """Add a shape to the slide"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=COLORS['white'], bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box to the slide"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=18, color=COLORS['white'], bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(6), font_name='Calibri'):
    """Add a paragraph to a text frame"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def create_title_slide(prs):
    """Slide 1: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Top and bottom accents
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=COLORS['gradient_blue'])
    add_shape(slide, Inches(0), Inches(7.0), Inches(13.33), Inches(0.08), fill_color=COLORS['gradient_purple'])
    
    # VS divider
    add_shape(slide, Inches(6.5), Inches(2.0), Inches(0.06), Inches(3.5), fill_color=COLORS['vs_red'])
    
    # Left title - Agent AI
    add_text_box(slide, Inches(1.5), Inches(2.2), Inches(4.5), Inches(1.2),
                 "Agent AI", font_size=52, color=COLORS['gradient_blue'], bold=True, alignment=PP_ALIGN.RIGHT)
    add_text_box(slide, Inches(1.5), Inches(3.4), Inches(4.5), Inches(0.6),
                 "代理式人工智能", font_size=24, color=COLORS['gradient_teal'], alignment=PP_ALIGN.RIGHT)
    
    # Right title - Agentic AI
    add_text_box(slide, Inches(7.3), Inches(2.2), Inches(4.5), Inches(1.2),
                 "Agentic AI", font_size=52, color=COLORS['gradient_purple'], bold=True, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, Inches(7.3), Inches(3.4), Inches(4.5), Inches(0.6),
                 "智能体人工智能", font_size=24, color=COLORS['accent'], alignment=PP_ALIGN.LEFT)
    
    # VS badge
    add_shape(slide, Inches(5.9), Inches(3.5), Inches(1.5), Inches(1.5), fill_color=COLORS['vs_red'], shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, Inches(5.9), Inches(3.8), Inches(1.5), Inches(0.8),
                 "VS", font_size=42, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
    
    # Subtitle
    add_text_box(slide, Inches(2), Inches(5.2), Inches(9.33), Inches(0.8),
                 "概念辨析 · 核心差异 · 技术演进 · 未来趋势", 
                 font_size=22, color=COLORS['text_gray'], alignment=PP_ALIGN.CENTER)
    
    # Bottom info
    add_text_box(slide, Inches(2), Inches(6.2), Inches(9.33), Inches(0.5),
                 "2026 | AI前沿深度解读", font_size=16, color=COLORS['text_gray'], alignment=PP_ALIGN.CENTER)

def create_concept_clarification_slide(prs):
    """Slide 2: Concept Clarification"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "概念澄清：一词之差，天壤之别", font_size=32, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=COLORS['gradient_blue'])
    
    # Agent AI - Left side
    add_shape(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.2), fill_color=COLORS['dark'])
    add_shape(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.08), fill_color=COLORS['gradient_blue'])
    
    add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.2), Inches(0.7),
                 "🔵 Agent AI（代理式AI）", font_size=30, color=COLORS['gradient_blue'], bold=True)

    agent_points = [
        "• 定义：以Agent为载体的AI系统",
        "• 重点：强调AI的存在形态",
        "• 能力：执行特定任务的代理程序",
        "• 交互：用户指令驱动",
        "• 典型：聊天机器人、虚拟助手",
    ]
    content = "\n".join(agent_points)
    add_text_box(slide, Inches(0.8), Inches(2.7), Inches(5.2), Inches(3.5),
                 content, font_size=18, color=COLORS['light'])

    add_shape(slide, Inches(0.8), Inches(5.5), Inches(5.2), Inches(0.03), fill_color=COLORS['gradient_blue'])
    add_text_box(slide, Inches(0.8), Inches(5.7), Inches(5.2), Inches(0.8),
                 "💡 本质：AI的\"载体\"与\"形态\"",
                 font_size=16, color=COLORS['gradient_teal'])
    
    # Agentic AI - Right side
    add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.2), fill_color=COLORS['dark'])
    add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.08), fill_color=COLORS['gradient_purple'])
    
    add_text_box(slide, Inches(7.3), Inches(1.8), Inches(5.2), Inches(0.7),
                 "🟣 Agentic AI（智能体AI）", font_size=30, color=COLORS['gradient_purple'], bold=True)

    agentic_points = [
        "• 定义：具有Agentic特性的AI系统",
        "• 重点：强调AI的行为与能力",
        "• 能力：自主规划、决策、执行",
        "• 交互：目标驱动，主动发起",
        "• 典型：AutoGPT、Devin",
    ]
    content = "\n".join(agentic_points)
    add_text_box(slide, Inches(7.3), Inches(2.7), Inches(5.2), Inches(3.5),
                 content, font_size=18, color=COLORS['light'])

    add_shape(slide, Inches(7.3), Inches(5.5), Inches(5.2), Inches(0.03), fill_color=COLORS['gradient_purple'])
    add_text_box(slide, Inches(7.3), Inches(5.7), Inches(5.2), Inches(0.8),
                 "💡 本质：AI的\"能力\"与\"自主性\"",
                 font_size=16, color=COLORS['accent'])

def create_definition_slide(prs):
    """Slide 3: What is Agentic AI"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "什么是 Agentic AI（智能体AI）", font_size=32, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=COLORS['gradient_purple'])
    
    # Definition box
    add_shape(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.5), fill_color=COLORS['dark'])
    add_shape(slide, Inches(0.8), Inches(1.6), Inches(0.08), Inches(1.5), fill_color=COLORS['gradient_purple'])
    
    add_text_box(slide, Inches(1.2), Inches(1.7), Inches(11), Inches(0.5),
                 "📖 定义", font_size=24, color=COLORS['gradient_purple'], bold=True)
    add_text_box(slide, Inches(1.2), Inches(2.2), Inches(11), Inches(0.8),
                 "Agentic AI 指能够自主感知环境、制定决策并采取行动以实现目标的智能系统。它不仅仅是响应指令，而是展现出主动性和目标导向行为。",
                 font_size=19, color=COLORS['light'])
    
    # Core characteristics - 3 columns
    characteristics = [
        ("🎯", "自主性", "Autonomy", [
            "• 无需人工干预",
            "• 独立制定计划",
            "• 自主执行任务",
            "• 自我纠错优化",
        ], COLORS['gradient_blue']),
        ("🚀", "主动性", "Proactiveness", [
            "• 主动发起行动",
            "• 预测用户需求",
            "• 提前规划路径",
            "• 动态调整策略",
        ], COLORS['secondary']),
        ("🤝", "交互性", "Social Ability", [
            "• 与其他AI协作",
            "• 人机自然对话",
            "• 多智能体协调",
            "• 环境感知反馈",
        ], COLORS['accent']),
    ]
    
    for idx, (icon, title, subtitle, points, color) in enumerate(characteristics):
        left = Inches(0.8 + idx * 4.1)
        
        # Card
        add_shape(slide, left, Inches(3.4), Inches(3.7), Inches(3.4), fill_color=COLORS['dark'])
        add_shape(slide, left, Inches(3.4), Inches(3.7), Inches(0.08), fill_color=color)
        
        # Icon & Title
        add_text_box(slide, left + Inches(0.2), Inches(3.6), Inches(0.8), Inches(0.6),
                     icon, font_size=32)
        add_text_box(slide, left + Inches(1.1), Inches(3.5), Inches(2.4), Inches(0.5),
                     title, font_size=24, color=color, bold=True)
        add_text_box(slide, left + Inches(1.1), Inches(4.0), Inches(2.4), Inches(0.4),
                     subtitle, font_size=14, color=COLORS['text_gray'])
        
        # Points
        content = "\n".join(points)
        add_text_box(slide, left + Inches(0.2), Inches(4.6), Inches(3.3), Inches(2),
                     content, font_size=16, color=COLORS['light'])

def create_comparison_slide(prs):
    """Slide 4: Detailed Comparison"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "核心维度对比", font_size=32, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=COLORS['accent'])
    
    # Table-like comparison
    comparisons = [
        ("对比维度", "Agent AI", "Agentic AI", COLORS['text_gray'], COLORS['gradient_blue'], COLORS['gradient_purple']),
        ("核心定位", "AI的载体形态", "AI的能力特征", COLORS['text_gray'], COLORS['gradient_blue'], COLORS['gradient_purple']),
        ("自主程度", "低-中（依赖指令）", "高（自主决策）", COLORS['text_gray'], COLORS['gradient_blue'], COLORS['gradient_purple']),
        ("任务处理", "执行预设任务", "自主规划执行", COLORS['text_gray'], COLORS['gradient_blue'], COLORS['gradient_purple']),
        ("交互方式", "被动响应", "主动发起", COLORS['text_gray'], COLORS['gradient_blue'], COLORS['gradient_purple']),
        ("决策能力", "有限/规则驱动", "强/推理驱动", COLORS['text_gray'], COLORS['gradient_blue'], COLORS['gradient_purple']),
        ("典型场景", "客服、助手、问答", "编程、研究、自动化", COLORS['text_gray'], COLORS['gradient_blue'], COLORS['gradient_purple']),
        ("发展阶段", "成熟/广泛应用", "快速演进/早期采用", COLORS['text_gray'], COLORS['gradient_blue'], COLORS['gradient_purple']),
    ]
    
    # Header row
    for idx, (text, _, _, bg_color, _, _) in enumerate(comparisons[:1]):
        left = Inches(0.8 + idx * 4.1)
        add_shape(slide, left, Inches(1.6), Inches(3.8), Inches(0.7), fill_color=COLORS['dark'])
        add_text_box(slide, left, Inches(1.65), Inches(3.8), Inches(0.6),
                     text, font_size=20, color=bg_color, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Column headers
    add_shape(slide, Inches(4.9), Inches(1.6), Inches(3.8), Inches(0.7), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(4.9), Inches(1.65), Inches(3.8), Inches(0.6),
                 "Agent AI 🔵", font_size=20, color=COLORS['gradient_blue'], bold=True, alignment=PP_ALIGN.CENTER)
    
    add_shape(slide, Inches(9.0), Inches(1.6), Inches(3.8), Inches(0.7), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(9.0), Inches(1.65), Inches(3.8), Inches(0.6),
                 "Agentic AI", font_size=20, color=COLORS['gradient_purple'], bold=True, alignment=PP_ALIGN.CENTER)
    
    # Data rows
    for row_idx, (dim, val1, val2, dim_color, _, _) in enumerate(comparisons[1:]):
        top = Inches(2.4 + row_idx * 0.65)
        bg = COLORS['dark'] if row_idx % 2 == 0 else COLORS['card_bg']
        
        # Dimension
        add_shape(slide, Inches(0.8), top, Inches(3.8), Inches(0.6), fill_color=bg)
        add_text_box(slide, Inches(0.9), top, Inches(3.6), Inches(0.55),
                     dim, font_size=16, color=dim_color, bold=True)
        
        # Value 1
        add_shape(slide, Inches(4.9), top, Inches(3.8), Inches(0.6), fill_color=bg)
        val1_color = COLORS['vs_red'] if val1 == "无" else COLORS['light']
        add_text_box(slide, Inches(5.0), top, Inches(3.6), Inches(0.55),
                     val1, font_size=15, color=val1_color)
        
        # Value 2
        add_shape(slide, Inches(9.0), top, Inches(3.8), Inches(0.6), fill_color=bg)
        add_text_box(slide, Inches(9.1), top, Inches(3.6), Inches(0.55),
                     val2, font_size=15, color=COLORS['light'])

def create_technical_architecture_slide(prs):
    """Slide 5: Technical Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "Agentic AI 技术架构", font_size=32, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=COLORS['gradient_blue'])
    
    # Architecture layers (bottom to top)
    layers = [
        ("🧠 认知层", "LLM + 知识库\n• 语言理解与推理\n• 上下文记忆管理\n• 领域知识调用", COLORS['gradient_blue']),
        ("⚙️ 规划层", "任务分解与编排\n• 目标分解规划\n• 工具选择调用\n• 执行路径优化", COLORS['gradient_purple']),
        ("🔧 工具层", "外部能力集成\n• API/数据库访问\n• 代码执行环境\n• 搜索引擎集成", COLORS['accent']),
        ("👁️ 感知层", "环境交互接口\n• 多模态输入处理\n• 实时反馈收集\n• 状态监测评估", COLORS['secondary']),
    ]
    
    for idx, (title, desc, color) in enumerate(layers):
        top = Inches(5.6 - idx * 1.2)  # Bottom to top
        
        # Layer card
        add_shape(slide, Inches(1.5), top, Inches(10.3), Inches(1.0), fill_color=COLORS['dark'])
        add_shape(slide, Inches(1.5), top, Inches(0.08), Inches(1.0), fill_color=color)
        
        # Title
        add_text_box(slide, Inches(1.8), top + Inches(0.1), Inches(2.5), Inches(0.5),
                     title, font_size=22, color=color, bold=True)
        
        # Description
        add_text_box(slide, Inches(4.3), top + Inches(0.1), Inches(7.3), Inches(0.8),
                     desc, font_size=15, color=COLORS['light'])
    
    # Architecture flow arrow
    add_text_box(slide, Inches(5.5), Inches(1.6), Inches(2.3), Inches(0.6),
                 "⬆ 数据流向上 ⬆", font_size=16, color=COLORS['text_gray'], alignment=PP_ALIGN.CENTER)

def create_use_cases_slide(prs):
    """Slide 6: Use Cases"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "Agentic AI 应用场景", font_size=32, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=COLORS['secondary'])
    
    # Use cases - 2x2 grid
    use_cases = [
        ("🤖", "自主编程助手", [
            "• 理解需求自动生成代码",
            "• 调试修复优化一体化",
            "• 代码审查与重构建议",
            "• 代表：Devin, Cursor Agent",
        ], COLORS['gradient_blue']),
        ("📊", "智能数据分析", [
            "• 自主探索数据发现模式",
            "• 生成可视化与报告",
            "• 统计建模与预测",
            "• 代表：Code Interpreter",
        ], COLORS['gradient_purple']),
        ("🌐", "自动化研究助理", [
            "• 文献检索与综述撰写",
            "• 实验设计与执行",
            "• 跨学科知识整合",
            "• 代表：Elicit, Consensus",
        ], COLORS['accent']),
        ("💼", "业务流程自动化", [
            "• 跨系统工作流编排",
            "• 客户服务自动化",
            "• 决策支持与执行",
            "• 代表：UiPath AI Agent",
        ], COLORS['secondary']),
    ]
    
    for idx, (icon, title, points, color) in enumerate(use_cases):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 6.1)
        top = Inches(1.6 + row * 2.8)
        
        # Card
        add_shape(slide, left, top, Inches(5.7), Inches(2.5), fill_color=COLORS['dark'])
        add_shape(slide, left, top, Inches(5.7), Inches(0.06), fill_color=color)
        
        # Icon & Title
        add_text_box(slide, left + Inches(0.3), top + Inches(0.2), Inches(0.8), Inches(0.6),
                     icon, font_size=32)
        add_text_box(slide, left + Inches(1.2), top + Inches(0.15), Inches(4.2), Inches(0.6),
                     title, font_size=24, color=color, bold=True)
        
        # Points
        content = "\n".join(points)
        add_text_box(slide, left + Inches(0.3), top + Inches(0.9), Inches(5.1), Inches(1.5),
                     content, font_size=16, color=COLORS['light'])

def create_evolution_slide(prs):
    """Slide 7: Evolution Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "AI演进时间线：从ChatGPT到Agentic AI", font_size=30, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=COLORS['gradient_teal'])
    
    # Timeline
    milestones = [
        ("2022\nQ4", "ChatGPT\n发布", "对话式AI\n引爆全球", COLORS['gradient_blue']),
        ("2023\nQ1", "GPT-4\n发布", "多模态能力\n初步智能", COLORS['gradient_purple']),
        ("2023\nQ2", "AutoGPT\n爆红", "AI智能体\n概念兴起", COLORS['accent']),
        ("2023\nQ4", "GPTs\n上线", "定制化\nAI助手", COLORS['secondary']),
        ("2024\nQ1", " Devin\n亮相", "自主编程\n智能体", COLORS['gradient_teal']),
        ("2025+", "Agentic AI\n时代", "多智能体\n协作生态", COLORS['gradient_blue']),
    ]
    
    for idx, (date, title, desc, color) in enumerate(milestones):
        left = Inches(0.8 + idx * 2.1)
        
        # Card
        add_shape(slide, left, Inches(1.8), Inches(1.8), Inches(4.5), fill_color=COLORS['dark'])
        add_shape(slide, left, Inches(1.8), Inches(1.8), Inches(0.06), fill_color=color)
        
        # Date
        add_text_box(slide, left, Inches(2.0), Inches(1.8), Inches(0.8),
                     date, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        
        # Divider
        add_shape(slide, Emu(left.emu + Inches(0.3).emu), Inches(2.9), Inches(1.2), Inches(0.02), fill_color=color)
        
        # Title
        add_text_box(slide, left + Inches(0.1), Inches(3.1), Inches(1.6), Inches(1.0),
                     title, font_size=20, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
        
        # Description
        add_text_box(slide, left + Inches(0.1), Inches(4.3), Inches(1.6), Inches(1.5),
                     desc, font_size=14, color=COLORS['text_gray'], alignment=PP_ALIGN.CENTER)
    
    # Arrow line
    add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.04), fill_color=COLORS['gradient_teal'])
    add_text_box(slide, Inches(4.5), Inches(6.3), Inches(4.3), Inches(0.5),
                 "AI自主性不断增强 ➡", font_size=18, color=COLORS['gradient_teal'], alignment=PP_ALIGN.CENTER)

def create_key_insights_slide(prs):
    """Slide 8: Key Insights"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "核心洞察与建议", font_size=32, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=COLORS['accent'])
    
    # Key insights
    insights = [
        ("💡", "洞察一：概念互补，非对立",
         "Agent AI强调AI的存在形态，Agentic AI强调AI的能力特征。\n两者不是对立关系，而是描述维度的不同。",
         COLORS['gradient_blue']),
        ("🎯", "洞察二：能力跃迁",
         "从Agent到Agentic，代表AI从\"被动执行\"到\"主动行动\"的范式转变。\n这是技术发展的必然趋势。",
         COLORS['gradient_purple']),
        ("🚀", "洞察三：融合发展趋势",
         "未来的AI系统将同时具备Agent形态和Agentic能力。\n两者的界限将逐渐模糊，走向统一。",
         COLORS['accent']),
        ("⚠️", "洞察四：风险意识",
         "Agentic AI的自主能力带来效率提升的同时，也带来可控性挑战。\n建立AI治理框架和监管机制刻不容缓。",
         COLORS['vs_red']),
    ]
    
    for idx, (icon, title, desc, color) in enumerate(insights):
        top = Inches(1.6 + idx * 1.35)
        
        # Card
        add_shape(slide, Inches(0.8), top, Inches(11.7), Inches(1.2), fill_color=COLORS['dark'])
        add_shape(slide, Inches(0.8), top, Inches(0.08), Inches(1.2), fill_color=color)
        
        # Icon & Title
        add_text_box(slide, Inches(1.1), top + Inches(0.15), Inches(0.6), Inches(0.5),
                     icon, font_size=24)
        add_text_box(slide, Inches(1.7), top + Inches(0.1), Inches(3), Inches(0.5),
                     title, font_size=20, color=color, bold=True)
        
        # Description
        add_text_box(slide, Inches(1.7), top + Inches(0.6), Inches(10.5), Inches(0.55),
                     desc, font_size=15, color=COLORS['light'])

def create_summary_slide(prs):
    """Slide 9: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header accent
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=COLORS['gradient_blue'])
    
    # Center summary
    add_shape(slide, Inches(3.5), Inches(1.5), Inches(6.3), Inches(4.5), fill_color=COLORS['dark'])
    add_shape(slide, Inches(3.5), Inches(1.5), Inches(6.3), Inches(0.08), fill_color=COLORS['gradient_purple'])
    add_shape(slide, Inches(3.5), Inches(5.92), Inches(6.3), Inches(0.08), fill_color=COLORS['gradient_blue'])
    
    add_text_box(slide, Inches(4), Inches(1.8), Inches(5.3), Inches(0.7),
                 "📌 总结", font_size=36, color=COLORS['gradient_purple'], bold=True, alignment=PP_ALIGN.CENTER)
    
    add_shape(slide, Inches(6.2), Inches(2.6), Inches(0.9), Inches(0.03), fill_color=COLORS['accent'])
    
    summary_points = [
        "🔵 Agent AI = AI的载体形态（在哪里）",
        "🟣 Agentic AI = AI的能力特征（能做什么）",
        "",
        "✓ Agent AI → 交互执行 · 任务导向",
        "✓ Agentic AI → 自主规划 · 目标驱动",
        "",
        "🔮 未来趋势：Agent + Agentic = 完整AI",
        "从被动工具 → 自主伙伴",
    ]
    
    content = "\n\n".join(summary_points)
    add_text_box(slide, Inches(4.2), Inches(2.9), Inches(4.9), Inches(2.8),
                 content, font_size=20, color=COLORS['light'], alignment=PP_ALIGN.CENTER)
    
    # Bottom tagline
    add_text_box(slide, Inches(2), Inches(6.3), Inches(9.33), Inches(0.6),
                 "理解差异，拥抱未来 · 智能体时代已来",
                 font_size=24, color=COLORS['gradient_blue'], bold=True, alignment=PP_ALIGN.CENTER)

def create_thank_you_slide(prs):
    """Slide 10: Thank You / Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Decorative elements
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=COLORS['gradient_blue'])
    add_shape(slide, Inches(0), Inches(7.0), Inches(13.33), Inches(0.08), fill_color=COLORS['gradient_purple'])
    
    # VS divider
    add_shape(slide, Inches(6.5), Inches(1.8), Inches(0.06), Inches(2.5), fill_color=COLORS['vs_red'])
    
    # Left
    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(4.5), Inches(0.8),
                 "Agent AI", font_size=42, color=COLORS['gradient_blue'], bold=True, alignment=PP_ALIGN.RIGHT)

    # Right
    add_text_box(slide, Inches(7.3), Inches(2.0), Inches(4.5), Inches(0.8),
                 "Agentic AI", font_size=42, color=COLORS['gradient_purple'], bold=True, alignment=PP_ALIGN.LEFT)

    # VS badge
    add_shape(slide, Inches(5.9), Inches(2.3), Inches(1.5), Inches(1.5), fill_color=COLORS['vs_red'], shape_type=MSO_SHAPE.OVAL)
    add_text_box(slide, Inches(5.9), Inches(2.6), Inches(1.5), Inches(0.8),
                 "VS", font_size=42, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)

    # Conclusion box
    add_shape(slide, Inches(4.0), Inches(4.5), Inches(5.3), Inches(1.0), fill_color=COLORS['dark'])
    add_shape(slide, Inches(4.0), Inches(4.5), Inches(5.3), Inches(0.06), fill_color=COLORS['secondary'])
    add_text_box(slide, Inches(4.0), Inches(4.6), Inches(5.3), Inches(0.4),
                 "✅ 两者相辅相成，共同定义AI未来", font_size=22, color=COLORS['secondary'], bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.0), Inches(5.0), Inches(5.3), Inches(0.4),
                 "Agent形态 + Agentic能力 = 完整智能", font_size=16, color=COLORS['light'], alignment=PP_ALIGN.CENTER)
    
    # Thank you
    add_text_box(slide, Inches(2), Inches(5.5), Inches(9.33), Inches(0.8),
                 "谢谢观看", font_size=48, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(2), Inches(6.3), Inches(9.33), Inches(0.5),
                 "Thank You | Q & A", font_size=22, color=COLORS['text_gray'], alignment=PP_ALIGN.CENTER)

def generate_ppt():
    """Generate the complete PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)
    
    # Create all slides
    create_title_slide(prs)
    create_concept_clarification_slide(prs)
    create_definition_slide(prs)
    create_comparison_slide(prs)
    create_technical_architecture_slide(prs)
    create_use_cases_slide(prs)
    create_evolution_slide(prs)
    create_key_insights_slide(prs)
    create_summary_slide(prs)
    create_thank_you_slide(prs)
    
    # Save
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Agent_AI_vs_Agentic_AI.pptx")
    prs.save(output_path)
    print(f"✅ PPT successfully generated: {output_path}")
    return output_path

if __name__ == "__main__":
    generate_ppt()
