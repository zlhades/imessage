#!/usr/bin/env python3
"""
AI Industrial Revolution - Professional PPT Generator
Creates a comprehensive presentation about the AI Industrial Revolution
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color scheme - Modern tech theme
COLORS = {
    'primary': RGBColor(0x1A, 0x73, 0xE8),      # Google Blue
    'secondary': RGBColor(0x34, 0xA8, 0x53),     # Green
    'accent': RGBColor(0xFF, 0x6D, 0x00),        # Orange
    'dark': RGBColor(0x20, 0x21, 0x24),          # Dark gray
    'darker': RGBColor(0x0D, 0x11, 0x17),        # Very dark
    'light': RGBColor(0xF8, 0xF9, 0xFA),         # Light gray
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'gradient_blue': RGBColor(0x42, 0x85, 0xF4),
    'gradient_purple': RGBColor(0x7C, 0x4D, 0xFF),
    'gradient_teal': RGBColor(0x00, 0xBC, 0xD4),
    'text_gray': RGBColor(0x5F, 0x63, 0x68),
    'card_bg': RGBColor(0x2D, 0x2F, 0x31),
}

def set_slide_bg(slide, color):
    """Set solid background color for a slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gradient_background(slide, color1, color2):
    """Add gradient-like background using shapes"""
    set_slide_bg(slide, color1)

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape_type=MSO_SHAPE.RECTANGLE):
    """Add a shape to the slide"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=COLORS['white'], bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box to the slide"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(text_frame, text, font_size=18, color=COLORS['white'], bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(6), font_name='Calibri'):
    """Add a paragraph to a text frame"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p

def create_title_slide(prs):
    """Slide 1: Title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_bg(slide, COLORS['darker'])
    
    # Add decorative elements
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=COLORS['primary'])
    add_shape(slide, Inches(0), Inches(7.0), Inches(13.33), Inches(0.08), fill_color=COLORS['primary'])
    
    # Side accent
    add_shape(slide, Inches(0.5), Inches(2.5), Inches(0.06), Inches(2), fill_color=COLORS['accent'])
    
    # Title
    add_text_box(slide, Inches(1.2), Inches(2.0), Inches(10), Inches(1.5),
                 "AI工业革命", font_size=54, color=COLORS['white'], bold=True)
    
    # Subtitle
    add_text_box(slide, Inches(1.2), Inches(3.5), Inches(10), Inches(1),
                 "Artificial Intelligence: The Fourth Industrial Revolution", 
                 font_size=28, color=COLORS['gradient_blue'], bold=False)
    
    # Decorative line
    add_shape(slide, Inches(1.2), Inches(4.6), Inches(3), Inches(0.04), fill_color=COLORS['accent'])
    
    # Info
    add_text_box(slide, Inches(1.2), Inches(5.0), Inches(8), Inches(0.8),
                 "从历史演变到未来趋势 | 重塑人类社会的智能浪潮", 
                 font_size=18, color=COLORS['text_gray'])
    
    # Bottom info
    add_text_box(slide, Inches(1.2), Inches(6.3), Inches(5), Inches(0.5),
                 "2026 | 深度研究报告", font_size=14, color=COLORS['text_gray'])

def create_timeline_slide(prs):
    """Slide 2: Industrial Revolution Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header bar
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "工业革命的历史演进", font_size=32, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=COLORS['primary'])
    
    # Timeline items
    timeline_data = [
        ("1.0", "蒸汽时代", "1760-1840", "蒸汽机、纺织机械", "机械化生产开始", COLORS['gradient_teal']),
        ("2.0", "电气时代", "1870-1914", "电力、流水线", "大规模生产成为可能", COLORS['gradient_blue']),
        ("3.0", "信息时代", "1960-2000", "计算机、互联网", "数字化革命兴起", COLORS['gradient_purple']),
        ("4.0", "智能时代", "2010-至今", "AI、物联网、大数据", "自主决策与智能化", COLORS['accent']),
    ]
    
    for i, (version, era, period, tech, desc, color) in enumerate(timeline_data):
        left = Inches(0.8 + i * 3.05)
        
        # Card background
        card = add_shape(slide, left, Inches(1.8), Inches(2.7), Inches(4.5), 
                        fill_color=COLORS['dark'])
        card.shadow.inherit = False
        
        # Top accent
        add_shape(slide, left, Inches(1.8), Inches(2.7), Inches(0.08), fill_color=color)
        
        # Version number
        add_text_box(slide, left, Inches(2.1), Inches(2.7), Inches(0.8),
                     f"工业{version}", font_size=28, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        
        # Era name
        add_text_box(slide, left, Inches(2.9), Inches(2.7), Inches(0.6),
                     era, font_size=24, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
        
        # Period
        add_text_box(slide, left, Inches(3.5), Inches(2.7), Inches(0.5),
                     period, font_size=16, color=COLORS['text_gray'], alignment=PP_ALIGN.CENTER)
        
        # Divider
        add_shape(slide, Emu(left.emu + Inches(0.5).emu), Inches(4.1), Inches(1.7), Inches(0.02), fill_color=color)
        
        # Technology
        tb = add_text_box(slide, left, Inches(4.3), Inches(2.7), Inches(1.5),
                         "核心技术", font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, Inches(4.7), Inches(2.7), Inches(1.2),
                     tech, font_size=16, color=COLORS['light'], alignment=PP_ALIGN.CENTER)
        
        # Description
        add_text_box(slide, left, Inches(5.5), Inches(2.7), Inches(0.8),
                     desc, font_size=13, color=COLORS['text_gray'], alignment=PP_ALIGN.CENTER)

def create_ai_core_tech_slide(prs):
    """Slide 3: AI Core Technologies"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(8), Inches(0.7),
                 "AI核心技术栈", font_size=32, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=COLORS['secondary'])
    
    # Tech cards in 2x2 grid
    tech_cards = [
        ("🧠", "机器学习", "Machine Learning", 
         "• 深度学习神经网络\n• 强化学习决策系统\n• 迁移学习跨域应用\n• 自监督学习范式", 
         COLORS['gradient_blue']),
        ("💬", "自然语言处理", "Natural Language Processing",
         "• 大语言模型 (LLM)\n• 语义理解与生成\n• 多语言实时翻译\n• 情感分析与洞察", 
         COLORS['gradient_purple']),
        ("👁️", "计算机视觉", "Computer Vision",
         "• 图像识别与分类\n• 目标检测与跟踪\n• 3D场景重建\n• 医学影像诊断", 
         COLORS['gradient_teal']),
        ("🤖", "智能机器人", "Robotics & Automation",
         "• 自主导航系统\n• 人机协作交互\n• 智能决策规划\n• 工业制造自动化", 
         COLORS['accent']),
    ]
    
    for idx, (icon, title, subtitle, points, color) in enumerate(tech_cards):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 6.1)
        top = Inches(1.6 + row * 2.9)
        
        # Card
        add_shape(slide, left, top, Inches(5.7), Inches(2.5), fill_color=COLORS['dark'])
        add_shape(slide, left, top, Inches(0.08), Inches(2.5), fill_color=color)
        
        # Icon & Title
        add_text_box(slide, left + Inches(0.3), top + Inches(0.2), Inches(1), Inches(0.6),
                     icon, font_size=32, color=color)
        add_text_box(slide, left + Inches(1.2), top + Inches(0.15), Inches(4), Inches(0.5),
                     title, font_size=24, color=COLORS['white'], bold=True)
        add_text_box(slide, left + Inches(1.2), top + Inches(0.6), Inches(4), Inches(0.4),
                     subtitle, font_size=14, color=color)
        
        # Points
        add_text_box(slide, left + Inches(0.3), top + Inches(1.2), Inches(5.2), Inches(1.2),
                     points, font_size=15, color=COLORS['light'])

def create_industry_applications_slide(prs):
    """Slide 4: Industry Applications"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "AI赋能千行百业", font_size=32, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=COLORS['accent'])
    
    # Application areas - 3 columns
    applications = [
        ("🏭 智能制造", [
            "• 质量检测自动化",
            "• 预测性维护系统",
            "• 供应链智能优化",
            "• 数字孪生工厂",
        ], COLORS['gradient_blue']),
        ("🏥 智慧医疗", [
            "• AI辅助疾病诊断",
            "• 药物研发加速",
            "• 个性化治疗方案",
            "• 医学影像分析",
        ], COLORS['secondary']),
        ("🚗 智能交通", [
            "• 自动驾驶技术",
            "• 智能交通调度",
            "• 物流路径优化",
            "• 车路协同系统",
        ], COLORS['gradient_purple']),
    ]
    
    for idx, (title, items, color) in enumerate(applications):
        left = Inches(0.8 + idx * 4.1)
        
        # Card
        add_shape(slide, left, Inches(1.6), Inches(3.7), Inches(4.8), fill_color=COLORS['dark'])
        add_shape(slide, left, Inches(1.6), Inches(3.7), Inches(0.1), fill_color=color)
        
        # Title
        add_text_box(slide, left + Inches(0.3), Inches(1.9), Inches(3.1), Inches(0.7),
                     title, font_size=24, color=color, bold=True)
        
        # Items
        content = "\n".join(items)
        add_text_box(slide, left + Inches(0.3), Inches(2.8), Inches(3.1), Inches(3.2),
                     content, font_size=17, color=COLORS['light'])
    
    # Bottom statistics
    stat_y = Inches(6.5)
    stats = [
        ("95%", "质检准确率"),
        ("70%", "研发效率提升"),
        ("40%", "物流成本降低"),
    ]
    
    for idx, (value, label) in enumerate(stats):
        left = Inches(1.5 + idx * 4)
        add_text_box(slide, left, stat_y, Inches(2), Inches(0.6),
                     value, font_size=36, color=COLORS['accent'], bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left, stat_y + Inches(0.6), Inches(2), Inches(0.4),
                     label, font_size=14, color=COLORS['text_gray'], alignment=PP_ALIGN.CENTER)

def create_economic_impact_slide(prs):
    """Slide 5: Economic Impact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "AI的经济影响力", font_size=32, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=COLORS['gradient_blue'])
    
    # Key metrics
    metrics = [
        ("$15.7万亿", "2030年全球AI\n经济贡献预测", COLORS['gradient_blue']),
        ("40%", "企业采用AI后\n效率平均提升", COLORS['secondary']),
        ("9700万", "AI创造的新就业\n岗位 (WEF 2025)", COLORS['accent']),
        ("$2.9万亿", "AI为企业带来的\n年度成本节约", COLORS['gradient_purple']),
    ]
    
    for idx, (value, desc, color) in enumerate(metrics):
        left = Inches(0.8 + idx * 3.1)
        
        # Card
        add_shape(slide, left, Inches(1.8), Inches(2.8), Inches(3.2), fill_color=COLORS['dark'])
        add_shape(slide, left, Inches(1.8), Inches(2.8), Inches(0.08), fill_color=color)

        # Value
        add_text_box(slide, left, Inches(2.2), Inches(2.8), Inches(1),
                     value, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, left + Inches(0.2), Inches(3.3), Inches(2.4), Inches(1.2),
                     desc, font_size=16, color=COLORS['light'], alignment=PP_ALIGN.CENTER)
    
    # Industry transformation section
    add_text_box(slide, Inches(0.8), Inches(5.3), Inches(5), Inches(0.6),
                 "产业变革趋势", font_size=24, color=COLORS['white'], bold=True)
    
    trends = [
        "1. 从自动化到自主化：AI系统从执行预设规则转向自主决策",
        "2. 从单一智能到群体智能：多AI系统协同工作的生态模式",
        "3. 从工具到伙伴：AI成为人类创造力的放大器",
        "4. 数据驱动决策：从经验判断到智能预测的范式转变",
    ]
    
    trend_box = add_text_box(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.2),
                             trends[0], font_size=16, color=COLORS['light'])
    tf = trend_box.text_frame
    for trend in trends[1:]:
        add_paragraph(tf, trend, font_size=16, color=COLORS['light'], space_before=Pt(8))

def create_challenges_slide(prs):
    """Slide 6: Challenges & Risks"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "AI革命的挑战与风险", font_size=32, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=RGBColor(0xFF, 0x44, 0x44))
    
    # Challenges in 2x2 grid
    challenges = [
        ("⚠️", "就业结构冲击", 
         "• 8500万岗位可能被自动化取代\n• 技能鸿沟扩大，再培训需求迫切\n• 需要全民终身学习体系",
         RGBColor(0xFF, 0x6B, 0x6B)),
        ("🔒", "隐私与数据安全",
         "• 海量数据采集引发隐私担忧\n• 数据泄露风险持续增加\n• 跨境数据流动监管难题",
         RGBColor(0xFF, 0xA7, 0x26)),
        ("⚖️", "伦理与偏见问题",
         "• 算法歧视加剧社会不平等\n• AI决策透明度与可解释性\n• 责任归属与问责机制缺失",
         RGBColor(0xAB, 0x47, 0xBC)),
        ("🌐", "治理与监管挑战",
         "• 国际AI军备竞赛风险\n• 技术垄断与数字鸿沟\n• 全球AI治理框架缺位",
         RGBColor(0x42, 0xA5, 0xF5)),
    ]
    
    for idx, (icon, title, points, color) in enumerate(challenges):
        row = idx // 2
        col = idx % 2
        left = Inches(0.8 + col * 6.1)
        top = Inches(1.6 + row * 2.8)
        
        # Card
        add_shape(slide, left, top, Inches(5.7), Inches(2.5), fill_color=COLORS['dark'])
        add_shape(slide, left, top, Inches(5.7), Inches(0.06), fill_color=color)
        
        # Icon & Title
        add_text_box(slide, left + Inches(0.3), top + Inches(0.3), Inches(0.8), Inches(0.5),
                     icon, font_size=28)
        add_text_box(slide, left + Inches(1.1), top + Inches(0.2), Inches(4.3), Inches(0.6),
                     title, font_size=22, color=color, bold=True)
        
        # Points
        add_text_box(slide, left + Inches(0.3), top + Inches(1.0), Inches(5.1), Inches(1.4),
                     points, font_size=14, color=COLORS['light'])

def create_future_vision_slide(prs):
    """Slide 7: Future Vision"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "未来展望：人机共生时代", font_size=32, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=COLORS['gradient_purple'])
    
    # Timeline to future
    future_milestones = [
        ("2025-2027", "AGI探索期", "通用人工智能初步探索\n多模态AI广泛应用", COLORS['gradient_blue']),
        ("2028-2032", "人机协作期", "脑机接口技术突破\nAI助手成为标配", COLORS['gradient_purple']),
        ("2033-2040", "智能融合期", "量子计算+AI结合\n自主AI系统出现", COLORS['accent']),
        ("2040+", "超级智能期", "超级智能可能性\n人类社会范式转变", COLORS['secondary']),
    ]
    
    for idx, (period, phase, desc, color) in enumerate(future_milestones):
        left = Inches(0.8 + idx * 3.1)
        
        # Card
        add_shape(slide, left, Inches(1.8), Inches(2.8), Inches(3.5), fill_color=COLORS['dark'])
        
        # Progress bar at top
        add_shape(slide, left, Inches(1.8), Inches(2.8), Inches(0.06), fill_color=color)

        # Period
        add_text_box(slide, left + Inches(0.2), Inches(2.1), Inches(2.4), Inches(0.5),
                     period, font_size=18, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        # Phase
        add_text_box(slide, left + Inches(0.2), Inches(2.7), Inches(2.4), Inches(0.7),
                     phase, font_size=22, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        add_text_box(slide, left + Inches(0.2), Inches(3.5), Inches(2.4), Inches(1.5),
                     desc, font_size=15, color=COLORS['light'], alignment=PP_ALIGN.CENTER)
    
    # Key insight box
    insight_left = Inches(0.8)
    insight_top = Inches(5.6)
    add_shape(slide, insight_left, insight_top, Inches(11.5), Inches(1.2), fill_color=COLORS['dark'])
    add_shape(slide, insight_left, insight_top, Inches(0.08), Inches(1.2), fill_color=COLORS['gradient_purple'])
    
    add_text_box(slide, insight_left + Inches(0.3), insight_top + Inches(0.15), Inches(10.8), Inches(0.9),
                 "💡 核心洞察：AI不会取代人类，但会使用AI的人将取代不会使用AI的人。未来的竞争力在于人机协作能力。",
                 font_size=18, color=COLORS['gradient_purple'], bold=True)

def create_action_plan_slide(prs):
    """Slide 8: Action Plan"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Header
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), fill_color=COLORS['dark'])
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.7),
                 "拥抱AI革命：行动指南", font_size=32, color=COLORS['white'], bold=True)
    add_shape(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), fill_color=COLORS['secondary'])
    
    # Action items for different stakeholders
    actions = [
        ("🎓 个人层面", [
            "1. 培养AI素养：学习基础AI知识",
            "2. 掌握AI工具：提升工作效率",
            "3. 发展软技能：创造力、批判性思维",
            "4. 终身学习：持续更新技能体系",
        ], COLORS['gradient_blue']),
        ("🏢 企业层面", [
            "1. 制定AI战略：明确数字化转型路径",
            "2. 投资人才培养：建设AI能力团队",
            "3. 构建数据基础：完善数据治理体系",
            "4. 敏捷试错：快速迭代AI应用场景",
        ], COLORS['secondary']),
        ("🏛️ 政府层面", [
            "1. 完善AI法规：建立治理框架",
            "2. 加大研发投入：支持基础研究",
            "3. 教育改革：重塑教育体系",
            "4. 社会保障：应对就业转型冲击",
        ], COLORS['accent']),
    ]
    
    for idx, (title, items, color) in enumerate(actions):
        left = Inches(0.8 + idx * 4.1)
        
        # Card
        add_shape(slide, left, Inches(1.6), Inches(3.7), Inches(4.5), fill_color=COLORS['dark'])
        add_shape(slide, left, Inches(1.6), Inches(3.7), Inches(0.1), fill_color=color)
        
        # Title
        add_text_box(slide, left + Inches(0.25), Inches(1.85), Inches(3.2), Inches(0.6),
                     title, font_size=22, color=color, bold=True)
        
        # Items
        content = "\n".join(items)
        add_text_box(slide, left + Inches(0.25), Inches(2.6), Inches(3.2), Inches(3.2),
                     content, font_size=15, color=COLORS['light'])
    
    # Bottom quote
    add_text_box(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.6),
                 "\"未来已来，只是分布不均。\" — William Gibson",
                 font_size=18, color=COLORS['text_gray'], alignment=PP_ALIGN.CENTER)

def create_thank_you_slide(prs):
    """Slide 9: Thank You / Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS['darker'])
    
    # Decorative elements
    add_shape(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.08), fill_color=COLORS['primary'])
    add_shape(slide, Inches(0), Inches(7.0), Inches(13.33), Inches(0.08), fill_color=COLORS['primary'])
    
    # Center content
    add_text_box(slide, Inches(2), Inches(2.0), Inches(9), Inches(1.5),
                 "谢谢观看", font_size=60, color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
    
    add_shape(slide, Inches(5.5), Inches(3.6), Inches(2.3), Inches(0.05), fill_color=COLORS['accent'])
    
    add_text_box(slide, Inches(2), Inches(4.0), Inches(9), Inches(0.8),
                 "Thank You", font_size=36, color=COLORS['gradient_blue'], alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(0.8),
                 "Q & A", font_size=28, color=COLORS['text_gray'], alignment=PP_ALIGN.CENTER)
    
    # Contact info placeholder
    add_text_box(slide, Inches(2), Inches(6.0), Inches(9), Inches(0.6),
                 "拥抱AI革命，共创智能未来",
                 font_size=20, color=COLORS['light'], alignment=PP_ALIGN.CENTER)

def generate_ppt():
    """Generate the complete PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)
    
    # Create all slides
    create_title_slide(prs)
    create_timeline_slide(prs)
    create_ai_core_tech_slide(prs)
    create_industry_applications_slide(prs)
    create_economic_impact_slide(prs)
    create_challenges_slide(prs)
    create_future_vision_slide(prs)
    create_action_plan_slide(prs)
    create_thank_you_slide(prs)
    
    # Save
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "AI_Industrial_Revolution.pptx")
    prs.save(output_path)
    print(f"✅ PPT successfully generated: {output_path}")
    return output_path

if __name__ == "__main__":
    generate_ppt()
